"""
Fallback document processors for when Docling service is unavailable.
"""

import logging
from pathlib import Path
from typing import Any, Dict, Optional
import asyncio
import time

import pipmaster as pm

logger = logging.getLogger(__name__)


class BasicDocumentProcessor:
    """Basic document processing using simple parsers."""
    
    @staticmethod 
    async def process_pdf(file_path: Path) -> str:
        """Process PDF using PyPDF2."""
        try:
            if not pm.is_installed("pypdf2"):
                pm.install("pypdf2")
            
            import PyPDF2
            
            content = ""
            with open(file_path, "rb") as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    content += page.extract_text() + "\n"
            
            return content.strip()
            
        except Exception as e:
            logger.error(f"PDF processing failed: {e}")
            return f"Error processing PDF: {str(e)}"
    
    @staticmethod
    async def process_docx(file_path: Path) -> str:
        """Process DOCX using python-docx."""
        try:
            if not pm.is_installed("python-docx"):
                try:
                    pm.install("python-docx")
                except Exception:
                    # Fallback to alternative package name
                    pm.install("docx")
            
            from docx import Document
            
            doc = Document(file_path)
            content = []
            
            for paragraph in doc.paragraphs:
                content.append(paragraph.text)
            
            return "\n".join(content).strip()
            
        except Exception as e:
            logger.error(f"DOCX processing failed: {e}")
            return f"Error processing DOCX: {str(e)}"
    
    @staticmethod
    async def process_pptx(file_path: Path) -> str:
        """Process PPTX using python-pptx."""
        try:
            if not pm.is_installed("python-pptx"):
                pm.install("python-pptx")
            
            from pptx import Presentation
            
            prs = Presentation(file_path)
            content = []
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content.append(shape.text)
            
            return "\n".join(content).strip()
            
        except Exception as e:
            logger.error(f"PPTX processing failed: {e}")
            return f"Error processing PPTX: {str(e)}"
    
    @staticmethod
    async def process_xlsx(file_path: Path) -> str:
        """Process XLSX using openpyxl."""
        try:
            if not pm.is_installed("openpyxl"):
                pm.install("openpyxl")
            
            from openpyxl import load_workbook
            
            workbook = load_workbook(file_path)
            content = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                content.append(f"Sheet: {sheet_name}")
                
                for row in sheet.iter_rows(values_only=True):
                    row_data = [str(cell) if cell is not None else "" for cell in row]
                    if any(row_data):  # Skip empty rows
                        content.append(" | ".join(row_data))
                
                content.append("")  # Empty line between sheets
            
            return "\n".join(content).strip()
            
        except Exception as e:
            logger.error(f"XLSX processing failed: {e}")
            return f"Error processing XLSX: {str(e)}"
    
    @staticmethod
    async def process_text(file_path: Path) -> str:
        """Process plain text files."""
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read().strip()
        except UnicodeDecodeError:
            # Try with different encoding
            try:
                with open(file_path, "r", encoding="latin1") as f:
                    return f.read().strip()
            except Exception as e:
                logger.error(f"Text processing failed: {e}")
                return f"Error reading text file: {str(e)}"
        except Exception as e:
            logger.error(f"Text processing failed: {e}")
            return f"Error reading text file: {str(e)}"
    
    async def process_document(self, file_path: Path, **kwargs) -> Dict[str, Any]:
        """Process document using appropriate basic parser."""
        start_time = time.time()
        file_path = Path(file_path)
        
        logger.info(f"Processing {file_path.name} with basic parser (Docling service unavailable)")
        
        try:
            # Determine processor based on file extension
            extension = file_path.suffix.lower()
            
            if extension == ".pdf":
                content = await self.process_pdf(file_path)
            elif extension == ".docx":
                content = await self.process_docx(file_path)
            elif extension == ".pptx":
                content = await self.process_pptx(file_path)
            elif extension == ".xlsx":
                content = await self.process_xlsx(file_path)
            elif extension in [".txt", ".md"]:
                content = await self.process_text(file_path)
            else:
                # Try to process as text for unknown extensions
                logger.warning(f"Unknown extension {extension}, trying as text")
                content = await self.process_text(file_path)
            
            processing_time = time.time() - start_time
            
            # Add basic metadata to match Docling service format
            metadata_section = self._generate_basic_metadata(
                file_path, 
                content, 
                processing_time,
                extension
            )
            
            # Add metadata to content if it was generated
            if metadata_section and kwargs.get("extract_metadata", True):
                full_content = metadata_section + "\n\n" + content
            else:
                full_content = content
            
            logger.info(f"Basic processing completed for {file_path.name} "
                       f"(processing_time: {processing_time:.2f}s, "
                       f"content_length: {len(full_content)})")
            
            return {
                "content": full_content,
                "metadata": {
                    "processing_time_seconds": processing_time,
                    "word_count": len(content.split()) if content else 0,
                    "character_count": len(content) if content else 0,
                    "processor": "basic_parser",
                    "file_extension": extension,
                    "docling_service_used": False,
                    "cache_hit": False
                },
                "success": True
            }
            
        except Exception as e:
            processing_time = time.time() - start_time
            error_msg = f"Basic processing failed for {file_path.name}: {str(e)}"
            logger.error(error_msg)
            
            return {
                "content": "",
                "metadata": {
                    "processing_time_seconds": processing_time,
                    "processor": "basic_parser",
                    "file_extension": file_path.suffix.lower(),
                    "docling_service_used": False,
                    "cache_hit": False,
                    "error": str(e)
                },
                "success": False,
                "error": error_msg
            }
    
    def _generate_basic_metadata(
        self, 
        file_path: Path, 
        content: str, 
        processing_time: float,
        extension: str
    ) -> str:
        """Generate basic metadata section for document."""
        from datetime import datetime, timezone
        
        lines = ["# Document Processing Metadata", ""]
        
        lines.append(f"- **Processed At**: {datetime.now(timezone.utc).isoformat()}")
        lines.append(f"- **Processing Time**: {processing_time:.2f} seconds")
        lines.append(f"- **Processor**: Basic Parser (Docling service unavailable)")
        lines.append(f"- **File Type**: {extension.upper()}")
        
        if content:
            word_count = len(content.split())
            char_count = len(content)
            lines.append(f"- **Word Count**: {word_count:,}")
            lines.append(f"- **Character Count**: {char_count:,}")
        
        lines.append("- **Features**: Basic text extraction only")
        lines.append("- **OCR**: Not available")
        lines.append("- **Table Recognition**: Not available")
        lines.append("- **Figure Extraction**: Not available")
        
        return "\n".join(lines)


class FallbackProcessor:
    """Manages fallback processing logic."""
    
    def __init__(self):
        self.basic_processor = BasicDocumentProcessor()
        self.fallback_enabled = True
        
    def set_fallback_enabled(self, enabled: bool):
        """Enable or disable fallback processing."""
        self.fallback_enabled = enabled
        logger.info(f"Fallback processing {'enabled' if enabled else 'disabled'}")
    
    async def process_document_with_fallback(
        self, 
        file_path: Path, 
        docling_service_available: bool = False,
        **kwargs
    ) -> Dict[str, Any]:
        """
        Process document with fallback logic.
        
        If Docling service is available, this should not be called.
        This is only for when the service is unavailable and fallback is enabled.
        """
        
        if not self.fallback_enabled:
            return {
                "content": "",
                "metadata": {
                    "processing_time_seconds": 0,
                    "processor": "none",
                    "docling_service_used": False,
                    "fallback_disabled": True
                },
                "success": False,
                "error": "Docling service unavailable and fallback disabled"
            }
        
        logger.info(f"Using fallback processor for {file_path.name}")
        
        # Use basic processor
        result = await self.basic_processor.process_document(file_path, **kwargs)
        
        # Add fallback indicator to metadata
        result["metadata"]["fallback_used"] = True
        result["metadata"]["docling_service_available"] = docling_service_available
        
        return result
    
    def get_supported_formats(self) -> list[str]:
        """Get list of formats supported by basic processors."""
        return [".pdf", ".docx", ".pptx", ".xlsx", ".txt", ".md"]
    
    def is_format_supported(self, file_path: Path) -> bool:
        """Check if file format is supported by basic processors."""
        extension = file_path.suffix.lower()
        return extension in self.get_supported_formats()


# Global fallback processor instance  
fallback_processor = FallbackProcessor()