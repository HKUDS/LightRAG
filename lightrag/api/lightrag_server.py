from fastapi import (
    FastAPI,
    HTTPException,
    File,
    UploadFile,
    BackgroundTasks,
)
import asyncio
import threading
import os
import json
import re
from fastapi.staticfiles import StaticFiles
import logging
import argparse
from typing import List, Any, Optional, Dict
from pydantic import BaseModel
from lightrag import LightRAG, QueryParam
from lightrag.types import GPTKeywordExtractionFormat
from lightrag.api import __api_version__
from lightrag.utils import EmbeddingFunc
from enum import Enum
from pathlib import Path
import shutil
import aiofiles
from ascii_colors import trace_exception, ASCIIColors
import sys
from fastapi import Depends, Security
from fastapi.security import APIKeyHeader
from fastapi.middleware.cors import CORSMiddleware
from contextlib import asynccontextmanager
from starlette.status import HTTP_403_FORBIDDEN
import pipmaster as pm
from dotenv import load_dotenv
import configparser
import traceback
from datetime import datetime

from lightrag.utils import logger
from .ollama_api import (
    OllamaAPI,
)
from .ollama_api import ollama_server_infos
from ..kg.postgres_impl import (
    PostgreSQLDB,
    PGKVStorage,
    PGVectorStorage,
    PGGraphStorage,
    PGDocStatusStorage,
)
from ..kg.oracle_impl import (
    OracleDB,
    OracleKVStorage,
    OracleVectorDBStorage,
    OracleGraphStorage,
)
from ..kg.tidb_impl import (
    TiDB,
    TiDBKVStorage,
    TiDBVectorDBStorage,
    TiDBGraphStorage,
)

# Load environment variables
load_dotenv(override=True)

# Initialize config parser
config = configparser.ConfigParser()
config.read("config.ini")


class DefaultRAGStorageConfig:
    KV_STORAGE = "JsonKVStorage"
    VECTOR_STORAGE = "NanoVectorDBStorage"
    GRAPH_STORAGE = "NetworkXStorage"
    DOC_STATUS_STORAGE = "JsonDocStatusStorage"


# Global progress tracker
scan_progress: Dict = {
    "is_scanning": False,
    "current_file": "",
    "indexed_count": 0,
    "total_files": 0,
    "progress": 0,
}

# Lock for thread-safe operations
progress_lock = threading.Lock()


def estimate_tokens(text: str) -> int:
    """Estimate the number of tokens in text
    Chinese characters: approximately 1.5 tokens per character
    English characters: approximately 0.25 tokens per character
    """
    # Use regex to match Chinese and non-Chinese characters separately
    chinese_chars = len(re.findall(r"[\u4e00-\u9fff]", text))
    non_chinese_chars = len(re.findall(r"[^\u4e00-\u9fff]", text))

    # Calculate estimated token count
    tokens = chinese_chars * 1.5 + non_chinese_chars * 0.25

    return int(tokens)


def get_default_host(binding_type: str) -> str:
    default_hosts = {
        "ollama": os.getenv("LLM_BINDING_HOST", "http://localhost:11434"),
        "lollms": os.getenv("LLM_BINDING_HOST", "http://localhost:9600"),
        "azure_openai": os.getenv("AZURE_OPENAI_ENDPOINT", "https://api.openai.com/v1"),
        "openai": os.getenv("LLM_BINDING_HOST", "https://api.openai.com/v1"),
    }
    return default_hosts.get(
        binding_type, os.getenv("LLM_BINDING_HOST", "http://localhost:11434")
    )  # fallback to ollama if unknown


def get_env_value(env_key: str, default: Any, value_type: type = str) -> Any:
    """
    Get value from environment variable with type conversion

    Args:
        env_key (str): Environment variable key
        default (Any): Default value if env variable is not set
        value_type (type): Type to convert the value to

    Returns:
        Any: Converted value from environment or default
    """
    value = os.getenv(env_key)
    if value is None:
        return default

    if isinstance(value_type, bool):
        return value.lower() in ("true", "1", "yes")
    try:
        return value_type(value)
    except ValueError:
        return default


def display_splash_screen(args: argparse.Namespace) -> None:
    """
    Display a colorful splash screen showing LightRAG server configuration

    Args:
        args: Parsed command line arguments
    """
    # Banner
    ASCIIColors.cyan(f"""
    ╔══════════════════════════════════════════════════════════════╗
    ║                   🚀 LightRAG Server v{__api_version__}                  ║
    ║          Fast, Lightweight RAG Server Implementation         ║
    ╚══════════════════════════════════════════════════════════════╝
    """)

    # Server Configuration
    ASCIIColors.magenta("\n📡 Server Configuration:")
    ASCIIColors.white("    ├─ Host: ", end="")
    ASCIIColors.yellow(f"{args.host}")
    ASCIIColors.white("    ├─ Port: ", end="")
    ASCIIColors.yellow(f"{args.port}")
    ASCIIColors.white("    ├─ CORS Origins: ", end="")
    ASCIIColors.yellow(f"{os.getenv('CORS_ORIGINS', '*')}")
    ASCIIColors.white("    ├─ SSL Enabled: ", end="")
    ASCIIColors.yellow(f"{args.ssl}")
    ASCIIColors.white("    └─ API Key: ", end="")
    ASCIIColors.yellow("Set" if args.key else "Not Set")
    if args.ssl:
        ASCIIColors.white("    ├─ SSL Cert: ", end="")
        ASCIIColors.yellow(f"{args.ssl_certfile}")
        ASCIIColors.white("    └─ SSL Key: ", end="")
        ASCIIColors.yellow(f"{args.ssl_keyfile}")

    # Directory Configuration
    ASCIIColors.magenta("\n📂 Directory Configuration:")
    ASCIIColors.white("    ├─ Working Directory: ", end="")
    ASCIIColors.yellow(f"{args.working_dir}")
    ASCIIColors.white("    └─ Input Directory: ", end="")
    ASCIIColors.yellow(f"{args.input_dir}")

    # LLM Configuration
    ASCIIColors.magenta("\n🤖 LLM Configuration:")
    ASCIIColors.white("    ├─ Binding: ", end="")
    ASCIIColors.yellow(f"{args.llm_binding}")
    ASCIIColors.white("    ├─ Host: ", end="")
    ASCIIColors.yellow(f"{args.llm_binding_host}")
    ASCIIColors.white("    └─ Model: ", end="")
    ASCIIColors.yellow(f"{args.llm_model}")

    # Embedding Configuration
    ASCIIColors.magenta("\n📊 Embedding Configuration:")
    ASCIIColors.white("    ├─ Binding: ", end="")
    ASCIIColors.yellow(f"{args.embedding_binding}")
    ASCIIColors.white("    ├─ Host: ", end="")
    ASCIIColors.yellow(f"{args.embedding_binding_host}")
    ASCIIColors.white("    ├─ Model: ", end="")
    ASCIIColors.yellow(f"{args.embedding_model}")
    ASCIIColors.white("    └─ Dimensions: ", end="")
    ASCIIColors.yellow(f"{args.embedding_dim}")

    # RAG Configuration
    ASCIIColors.magenta("\n⚙️ RAG Configuration:")
    ASCIIColors.white("    ├─ Max Async Operations: ", end="")
    ASCIIColors.yellow(f"{args.max_async}")
    ASCIIColors.white("    ├─ Max Tokens: ", end="")
    ASCIIColors.yellow(f"{args.max_tokens}")
    ASCIIColors.white("    ├─ Max Embed Tokens: ", end="")
    ASCIIColors.yellow(f"{args.max_embed_tokens}")
    ASCIIColors.white("    ├─ Chunk Size: ", end="")
    ASCIIColors.yellow(f"{args.chunk_size}")
    ASCIIColors.white("    ├─ Chunk Overlap Size: ", end="")
    ASCIIColors.yellow(f"{args.chunk_overlap_size}")
    ASCIIColors.white("    ├─ History Turns: ", end="")
    ASCIIColors.yellow(f"{args.history_turns}")
    ASCIIColors.white("    ├─ Cosine Threshold: ", end="")
    ASCIIColors.yellow(f"{args.cosine_threshold}")
    ASCIIColors.white("    └─ Top-K: ", end="")
    ASCIIColors.yellow(f"{args.top_k}")

    # System Configuration
    ASCIIColors.magenta("\n💾 Storage Configuration:")
    ASCIIColors.white("    ├─ KV Storage: ", end="")
    ASCIIColors.yellow(f"{args.kv_storage}")
    ASCIIColors.white("    ├─ Vector Storage: ", end="")
    ASCIIColors.yellow(f"{args.vector_storage}")
    ASCIIColors.white("    ├─ Graph Storage: ", end="")
    ASCIIColors.yellow(f"{args.graph_storage}")
    ASCIIColors.white("    └─ Document Status Storage: ", end="")
    ASCIIColors.yellow(f"{args.doc_status_storage}")

    ASCIIColors.magenta("\n🛠️ System Configuration:")
    ASCIIColors.white("    ├─ Ollama Emulating Model: ", end="")
    ASCIIColors.yellow(f"{ollama_server_infos.LIGHTRAG_MODEL}")
    ASCIIColors.white("    ├─ Log Level: ", end="")
    ASCIIColors.yellow(f"{args.log_level}")
    ASCIIColors.white("    └─ Timeout: ", end="")
    ASCIIColors.yellow(f"{args.timeout if args.timeout else 'None (infinite)'}")

    # Server Status
    ASCIIColors.green("\n✨ Server starting up...\n")

    # Server Access Information
    protocol = "https" if args.ssl else "http"
    if args.host == "0.0.0.0":
        ASCIIColors.magenta("\n🌐 Server Access Information:")
        ASCIIColors.white("    ├─ Local Access: ", end="")
        ASCIIColors.yellow(f"{protocol}://localhost:{args.port}")
        ASCIIColors.white("    ├─ Remote Access: ", end="")
        ASCIIColors.yellow(f"{protocol}://<your-ip-address>:{args.port}")
        ASCIIColors.white("    ├─ API Documentation (local): ", end="")
        ASCIIColors.yellow(f"{protocol}://localhost:{args.port}/docs")
        ASCIIColors.white("    ├─ Alternative Documentation (local): ", end="")
        ASCIIColors.yellow(f"{protocol}://localhost:{args.port}/redoc")
        ASCIIColors.white("    ├─ WebUI (local): ", end="")
        ASCIIColors.yellow(f"{protocol}://localhost:{args.port}/webui")
        ASCIIColors.white("    └─ Graph Viewer (local): ", end="")
        ASCIIColors.yellow(f"{protocol}://localhost:{args.port}/graph-viewer")

        ASCIIColors.yellow("\n📝 Note:")
        ASCIIColors.white("""    Since the server is running on 0.0.0.0:
    - Use 'localhost' or '127.0.0.1' for local access
    - Use your machine's IP address for remote access
    - To find your IP address:
      • Windows: Run 'ipconfig' in terminal
      • Linux/Mac: Run 'ifconfig' or 'ip addr' in terminal
    """)
    else:
        base_url = f"{protocol}://{args.host}:{args.port}"
        ASCIIColors.magenta("\n🌐 Server Access Information:")
        ASCIIColors.white("    ├─ Base URL: ", end="")
        ASCIIColors.yellow(f"{base_url}")
        ASCIIColors.white("    ├─ API Documentation: ", end="")
        ASCIIColors.yellow(f"{base_url}/docs")
        ASCIIColors.white("    └─ Alternative Documentation: ", end="")
        ASCIIColors.yellow(f"{base_url}/redoc")

    # Usage Examples
    ASCIIColors.magenta("\n📚 Quick Start Guide:")
    ASCIIColors.cyan("""
    1. Access the Swagger UI:
       Open your browser and navigate to the API documentation URL above

    2. API Authentication:""")
    if args.key:
        ASCIIColors.cyan("""       Add the following header to your requests:
       X-API-Key: <your-api-key>
    """)
    else:
        ASCIIColors.cyan("       No authentication required\n")

    ASCIIColors.cyan("""    3. Basic Operations:
       - POST /upload_document: Upload new documents to RAG
       - POST /query: Query your document collection
       - GET /collections: List available collections

    4. Monitor the server:
       - Check server logs for detailed operation information
       - Use healthcheck endpoint: GET /health
    """)

    # Security Notice
    if args.key:
        ASCIIColors.yellow("\n⚠️  Security Notice:")
        ASCIIColors.white("""    API Key authentication is enabled.
    Make sure to include the X-API-Key header in all your requests.
    """)

    ASCIIColors.green("Server is ready to accept connections! 🚀\n")

    # Ensure splash output flush to system log
    sys.stdout.flush()


def parse_args() -> argparse.Namespace:
    """
    Parse command line arguments with environment variable fallback

    Returns:
        argparse.Namespace: Parsed arguments
    """

    parser = argparse.ArgumentParser(
        description="LightRAG FastAPI Server with separate working and input directories"
    )

    parser.add_argument(
        "--kv-storage",
        default=get_env_value(
            "LIGHTRAG_KV_STORAGE", DefaultRAGStorageConfig.KV_STORAGE
        ),
        help=f"KV存储实现 (default: {DefaultRAGStorageConfig.KV_STORAGE})",
    )
    parser.add_argument(
        "--doc-status-storage",
        default=get_env_value(
            "LIGHTRAG_DOC_STATUS_STORAGE", DefaultRAGStorageConfig.DOC_STATUS_STORAGE
        ),
        help=f"文档状态存储实现 (default: {DefaultRAGStorageConfig.DOC_STATUS_STORAGE})",
    )
    parser.add_argument(
        "--graph-storage",
        default=get_env_value(
            "LIGHTRAG_GRAPH_STORAGE", DefaultRAGStorageConfig.GRAPH_STORAGE
        ),
        help=f"图存储实现 (default: {DefaultRAGStorageConfig.GRAPH_STORAGE})",
    )
    parser.add_argument(
        "--vector-storage",
        default=get_env_value(
            "LIGHTRAG_VECTOR_STORAGE", DefaultRAGStorageConfig.VECTOR_STORAGE
        ),
        help=f"向量存储实现 (default: {DefaultRAGStorageConfig.VECTOR_STORAGE})",
    )

    # Bindings configuration
    parser.add_argument(
        "--llm-binding",
        default=get_env_value("LLM_BINDING", "ollama"),
        help="LLM binding to be used. Supported: lollms, ollama, openai (default: from env or ollama)",
    )
    parser.add_argument(
        "--embedding-binding",
        default=get_env_value("EMBEDDING_BINDING", "ollama"),
        help="Embedding binding to be used. Supported: lollms, ollama, openai (default: from env or ollama)",
    )

    # Server configuration
    parser.add_argument(
        "--host",
        default=get_env_value("HOST", "0.0.0.0"),
        help="Server host (default: from env or 0.0.0.0)",
    )
    parser.add_argument(
        "--port",
        type=int,
        default=get_env_value("PORT", 9621, int),
        help="Server port (default: from env or 9621)",
    )

    # Directory configuration
    parser.add_argument(
        "--working-dir",
        default=get_env_value("WORKING_DIR", "./rag_storage"),
        help="Working directory for RAG storage (default: from env or ./rag_storage)",
    )
    parser.add_argument(
        "--input-dir",
        default=get_env_value("INPUT_DIR", "./inputs"),
        help="Directory containing input documents (default: from env or ./inputs)",
    )

    # LLM Model configuration
    parser.add_argument(
        "--llm-binding-host",
        default=get_env_value("LLM_BINDING_HOST", None),
        help="LLM server host URL. If not provided, defaults based on llm-binding:\n"
        + "- ollama: http://localhost:11434\n"
        + "- lollms: http://localhost:9600\n"
        + "- openai: https://api.openai.com/v1",
    )

    default_llm_api_key = get_env_value("LLM_BINDING_API_KEY", None)

    parser.add_argument(
        "--llm-binding-api-key",
        default=default_llm_api_key,
        help="llm server API key (default: from env or empty string)",
    )

    parser.add_argument(
        "--llm-model",
        default=get_env_value("LLM_MODEL", "mistral-nemo:latest"),
        help="LLM model name (default: from env or mistral-nemo:latest)",
    )

    # Embedding model configuration
    parser.add_argument(
        "--embedding-binding-host",
        default=get_env_value("EMBEDDING_BINDING_HOST", None),
        help="Embedding server host URL. If not provided, defaults based on embedding-binding:\n"
        + "- ollama: http://localhost:11434\n"
        + "- lollms: http://localhost:9600\n"
        + "- openai: https://api.openai.com/v1",
    )

    default_embedding_api_key = get_env_value("EMBEDDING_BINDING_API_KEY", "")
    parser.add_argument(
        "--embedding-binding-api-key",
        default=default_embedding_api_key,
        help="embedding server API key (default: from env or empty string)",
    )

    parser.add_argument(
        "--embedding-model",
        default=get_env_value("EMBEDDING_MODEL", "bge-m3:latest"),
        help="Embedding model name (default: from env or bge-m3:latest)",
    )

    parser.add_argument(
        "--chunk_size",
        default=get_env_value("CHUNK_SIZE", 1200),
        help="chunk chunk size default 1200",
    )

    parser.add_argument(
        "--chunk_overlap_size",
        default=get_env_value("CHUNK_OVERLAP_SIZE", 100),
        help="chunk overlap size default 100",
    )

    def timeout_type(value):
        if value is None or value == "None":
            return None
        return int(value)

    parser.add_argument(
        "--timeout",
        default=get_env_value("TIMEOUT", None, timeout_type),
        type=timeout_type,
        help="Timeout in seconds (useful when using slow AI). Use None for infinite timeout",
    )

    # RAG configuration
    parser.add_argument(
        "--max-async",
        type=int,
        default=get_env_value("MAX_ASYNC", 4, int),
        help="Maximum async operations (default: from env or 4)",
    )
    parser.add_argument(
        "--max-tokens",
        type=int,
        default=get_env_value("MAX_TOKENS", 32768, int),
        help="Maximum token size (default: from env or 32768)",
    )
    parser.add_argument(
        "--embedding-dim",
        type=int,
        default=get_env_value("EMBEDDING_DIM", 1024, int),
        help="Embedding dimensions (default: from env or 1024)",
    )
    parser.add_argument(
        "--max-embed-tokens",
        type=int,
        default=get_env_value("MAX_EMBED_TOKENS", 8192, int),
        help="Maximum embedding token size (default: from env or 8192)",
    )

    # Logging configuration
    parser.add_argument(
        "--log-level",
        default=get_env_value("LOG_LEVEL", "INFO"),
        choices=["DEBUG", "INFO", "WARNING", "ERROR", "CRITICAL"],
        help="Logging level (default: from env or INFO)",
    )

    parser.add_argument(
        "--key",
        type=str,
        default=get_env_value("LIGHTRAG_API_KEY", None),
        help="API key for authentication. This protects lightrag server against unauthorized access",
    )

    # Optional https parameters
    parser.add_argument(
        "--ssl",
        action="store_true",
        default=get_env_value("SSL", False, bool),
        help="Enable HTTPS (default: from env or False)",
    )
    parser.add_argument(
        "--ssl-certfile",
        default=get_env_value("SSL_CERTFILE", None),
        help="Path to SSL certificate file (required if --ssl is enabled)",
    )
    parser.add_argument(
        "--ssl-keyfile",
        default=get_env_value("SSL_KEYFILE", None),
        help="Path to SSL private key file (required if --ssl is enabled)",
    )
    parser.add_argument(
        "--auto-scan-at-startup",
        action="store_true",
        default=False,
        help="Enable automatic scanning when the program starts",
    )

    parser.add_argument(
        "--history-turns",
        type=int,
        default=get_env_value("HISTORY_TURNS", 3, int),
        help="Number of conversation history turns to include (default: from env or 3)",
    )

    # Search parameters
    parser.add_argument(
        "--top-k",
        type=int,
        default=get_env_value("TOP_K", 60, int),
        help="Number of most similar results to return (default: from env or 60)",
    )
    parser.add_argument(
        "--cosine-threshold",
        type=float,
        default=get_env_value("COSINE_THRESHOLD", 0.2, float),
        help="Cosine similarity threshold (default: from env or 0.4)",
    )

    # Ollama model name
    parser.add_argument(
        "--simulated-model-name",
        type=str,
        default=get_env_value(
            "SIMULATED_MODEL_NAME", ollama_server_infos.LIGHTRAG_MODEL
        ),
        help="Number of conversation history turns to include (default: from env or 3)",
    )

    # Namespace
    parser.add_argument(
        "--namespace-prefix",
        type=str,
        default=get_env_value("NAMESPACE_PREFIX", ""),
        help="Prefix of the namespace",
    )

    args = parser.parse_args()

    # conver relative path to absolute path
    args.working_dir = os.path.abspath(args.working_dir)
    args.input_dir = os.path.abspath(args.input_dir)

    ollama_server_infos.LIGHTRAG_MODEL = args.simulated_model_name

    return args


class DocumentManager:
    """Handles document operations and tracking"""

    def __init__(
        self,
        input_dir: str,
        supported_extensions: tuple = (
            ".txt",
            ".md",
            ".pdf",
            ".docx",
            ".pptx",
            ".xlsx",
        ),
    ):
        self.input_dir = Path(input_dir)
        self.supported_extensions = supported_extensions
        self.indexed_files = set()

        # Create input directory if it doesn't exist
        self.input_dir.mkdir(parents=True, exist_ok=True)

    def scan_directory_for_new_files(self) -> List[Path]:
        """Scan input directory for new files"""
        new_files = []
        for ext in self.supported_extensions:
            logger.info(f"Scanning for {ext} files in {self.input_dir}")
            for file_path in self.input_dir.rglob(f"*{ext}"):
                if file_path not in self.indexed_files:
                    new_files.append(file_path)
        return new_files

    def scan_directory(self) -> List[Path]:
        """Scan input directory for new files"""
        new_files = []
        for ext in self.supported_extensions:
            for file_path in self.input_dir.rglob(f"*{ext}"):
                new_files.append(file_path)
        return new_files

    def mark_as_indexed(self, file_path: Path):
        """Mark a file as indexed"""
        self.indexed_files.add(file_path)

    def is_supported_file(self, filename: str) -> bool:
        """Check if file type is supported"""
        return any(filename.lower().endswith(ext) for ext in self.supported_extensions)


# LightRAG query mode
class SearchMode(str, Enum):
    naive = "naive"
    local = "local"
    global_ = "global"
    hybrid = "hybrid"
    mix = "mix"


class QueryRequest(BaseModel):
    query: str

    """Specifies the retrieval mode"""
    mode: SearchMode = SearchMode.hybrid

    """If True, enables streaming output for real-time responses."""
    stream: Optional[bool] = None

    """If True, only returns the retrieved context without generating a response."""
    only_need_context: Optional[bool] = None

    """If True, only returns the generated prompt without producing a response."""
    only_need_prompt: Optional[bool] = None

    """Defines the response format. Examples: 'Multiple Paragraphs', 'Single Paragraph', 'Bullet Points'."""
    response_type: Optional[str] = None

    """Number of top items to retrieve. Represents entities in 'local' mode and relationships in 'global' mode."""
    top_k: Optional[int] = None

    """Maximum number of tokens allowed for each retrieved text chunk."""
    max_token_for_text_unit: Optional[int] = None

    """Maximum number of tokens allocated for relationship descriptions in global retrieval."""
    max_token_for_global_context: Optional[int] = None

    """Maximum number of tokens allocated for entity descriptions in local retrieval."""
    max_token_for_local_context: Optional[int] = None

    """List of high-level keywords to prioritize in retrieval."""
    hl_keywords: Optional[List[str]] = None

    """List of low-level keywords to refine retrieval focus."""
    ll_keywords: Optional[List[str]] = None

    """Stores past conversation history to maintain context.
    Format: [{"role": "user/assistant", "content": "message"}].
    """
    conversation_history: Optional[List[dict[str, Any]]] = None

    """Number of complete conversation turns (user-assistant pairs) to consider in the response context."""
    history_turns: Optional[int] = None


class QueryResponse(BaseModel):
    response: str


class InsertTextRequest(BaseModel):
    text: str


class InsertResponse(BaseModel):
    status: str
    message: str


def QueryRequestToQueryParams(request: QueryRequest):
    param = QueryParam(mode=request.mode, stream=request.stream)
    if request.only_need_context is not None:
        param.only_need_context = request.only_need_context
    if request.only_need_prompt is not None:
        param.only_need_prompt = request.only_need_prompt
    if request.response_type is not None:
        param.response_type = request.response_type
    if request.top_k is not None:
        param.top_k = request.top_k
    if request.max_token_for_text_unit is not None:
        param.max_token_for_text_unit = request.max_token_for_text_unit
    if request.max_token_for_global_context is not None:
        param.max_token_for_global_context = request.max_token_for_global_context
    if request.max_token_for_local_context is not None:
        param.max_token_for_local_context = request.max_token_for_local_context
    if request.hl_keywords is not None:
        param.hl_keywords = request.hl_keywords
    if request.ll_keywords is not None:
        param.ll_keywords = request.ll_keywords
    if request.conversation_history is not None:
        param.conversation_history = request.conversation_history
    if request.history_turns is not None:
        param.history_turns = request.history_turns
    return param


def get_api_key_dependency(api_key: Optional[str]):
    if not api_key:
        # If no API key is configured, return a dummy dependency that always succeeds
        async def no_auth():
            return None

        return no_auth

    # If API key is configured, use proper authentication
    api_key_header = APIKeyHeader(name="X-API-Key", auto_error=False)

    async def api_key_auth(
        api_key_header_value: Optional[str] = Security(api_key_header),
    ):
        if not api_key_header_value:
            raise HTTPException(
                status_code=HTTP_403_FORBIDDEN, detail="API Key required"
            )
        if api_key_header_value != api_key:
            raise HTTPException(
                status_code=HTTP_403_FORBIDDEN, detail="Invalid API Key"
            )
        return api_key_header_value

    return api_key_auth


# Global configuration
global_top_k = 60  # default value
temp_prefix = "__tmp_"  # prefix for temporary files


def create_app(args):
    global global_top_k
    global_top_k = args.top_k  # save top_k from args

    # Verify that bindings are correctly setup
    if args.llm_binding not in [
        "lollms",
        "ollama",
        "openai",
        "openai-ollama",
        "azure_openai",
    ]:
        raise Exception("llm binding not supported")

    if args.embedding_binding not in ["lollms", "ollama", "openai", "azure_openai"]:
        raise Exception("embedding binding not supported")

    # Set default hosts if not provided
    if args.llm_binding_host is None:
        args.llm_binding_host = get_default_host(args.llm_binding)

    if args.embedding_binding_host is None:
        args.embedding_binding_host = get_default_host(args.embedding_binding)

    # Add SSL validation
    if args.ssl:
        if not args.ssl_certfile or not args.ssl_keyfile:
            raise Exception(
                "SSL certificate and key files must be provided when SSL is enabled"
            )
        if not os.path.exists(args.ssl_certfile):
            raise Exception(f"SSL certificate file not found: {args.ssl_certfile}")
        if not os.path.exists(args.ssl_keyfile):
            raise Exception(f"SSL key file not found: {args.ssl_keyfile}")

    # Setup logging
    logging.basicConfig(
        format="%(levelname)s:%(message)s", level=getattr(logging, args.log_level)
    )

    # Check if API key is provided either through env var or args
    api_key = os.getenv("LIGHTRAG_API_KEY") or args.key

    # Initialize document manager
    doc_manager = DocumentManager(args.input_dir)

    @asynccontextmanager
    async def lifespan(app: FastAPI):
        """Lifespan context manager for startup and shutdown events"""
        # Initialize database connections
        postgres_db = None
        oracle_db = None
        tidb_db = None
        # Store background tasks
        app.state.background_tasks = set()

        try:
            # Check if PostgreSQL is needed
            if any(
                isinstance(
                    storage_instance,
                    (PGKVStorage, PGVectorStorage, PGGraphStorage, PGDocStatusStorage),
                )
                for _, storage_instance in storage_instances
            ):
                postgres_db = PostgreSQLDB(_get_postgres_config())
                await postgres_db.initdb()
                await postgres_db.check_tables()
                for storage_name, storage_instance in storage_instances:
                    if isinstance(
                        storage_instance,
                        (
                            PGKVStorage,
                            PGVectorStorage,
                            PGGraphStorage,
                            PGDocStatusStorage,
                        ),
                    ):
                        storage_instance.db = postgres_db
                        logger.info(f"Injected postgres_db to {storage_name}")

            # Check if Oracle is needed
            if any(
                isinstance(
                    storage_instance,
                    (OracleKVStorage, OracleVectorDBStorage, OracleGraphStorage),
                )
                for _, storage_instance in storage_instances
            ):
                oracle_db = OracleDB(_get_oracle_config())
                await oracle_db.check_tables()
                for storage_name, storage_instance in storage_instances:
                    if isinstance(
                        storage_instance,
                        (OracleKVStorage, OracleVectorDBStorage, OracleGraphStorage),
                    ):
                        storage_instance.db = oracle_db
                        logger.info(f"Injected oracle_db to {storage_name}")

            # Check if TiDB is needed
            if any(
                isinstance(
                    storage_instance,
                    (TiDBKVStorage, TiDBVectorDBStorage, TiDBGraphStorage),
                )
                for _, storage_instance in storage_instances
            ):
                tidb_db = TiDB(_get_tidb_config())
                await tidb_db.check_tables()
                for storage_name, storage_instance in storage_instances:
                    if isinstance(
                        storage_instance,
                        (TiDBKVStorage, TiDBVectorDBStorage, TiDBGraphStorage),
                    ):
                        storage_instance.db = tidb_db
                        logger.info(f"Injected tidb_db to {storage_name}")

            # Auto scan documents if enabled
            if args.auto_scan_at_startup:
                # Start scanning in background
                with progress_lock:
                    if not scan_progress["is_scanning"]:
                        scan_progress["is_scanning"] = True
                        scan_progress["indexed_count"] = 0
                        scan_progress["progress"] = 0
                        # Create background task
                        task = asyncio.create_task(run_scanning_process())
                        app.state.background_tasks.add(task)
                        task.add_done_callback(app.state.background_tasks.discard)
                        ASCIIColors.info(
                            f"Started background scanning of documents from {args.input_dir}"
                        )
                    else:
                        ASCIIColors.info(
                            "Skip document scanning(anohter scanning is active)"
                        )

            yield

        finally:
            # Cleanup database connections
            if postgres_db and hasattr(postgres_db, "pool"):
                await postgres_db.pool.close()
                logger.info("Closed PostgreSQL connection pool")

            if oracle_db and hasattr(oracle_db, "pool"):
                await oracle_db.pool.close()
                logger.info("Closed Oracle connection pool")

            if tidb_db and hasattr(tidb_db, "pool"):
                await tidb_db.pool.close()
                logger.info("Closed TiDB connection pool")

    # Initialize FastAPI
    app = FastAPI(
        title="LightRAG API",
        description="API for querying text using LightRAG with separate storage and input directories"
        + "(With authentication)"
        if api_key
        else "",
        version=__api_version__,
        openapi_tags=[{"name": "api"}],
        lifespan=lifespan,
    )

    def get_cors_origins():
        """Get allowed origins from environment variable
        Returns a list of allowed origins, defaults to ["*"] if not set
        """
        origins_str = os.getenv("CORS_ORIGINS", "*")
        if origins_str == "*":
            return ["*"]
        return [origin.strip() for origin in origins_str.split(",")]

    # Add CORS middleware
    app.add_middleware(
        CORSMiddleware,
        allow_origins=get_cors_origins(),
        allow_credentials=True,
        allow_methods=["*"],
        allow_headers=["*"],
    )

    # Database configuration functions
    def _get_postgres_config():
        return {
            "host": os.environ.get(
                "POSTGRES_HOST",
                config.get("postgres", "host", fallback="localhost"),
            ),
            "port": os.environ.get(
                "POSTGRES_PORT", config.get("postgres", "port", fallback=5432)
            ),
            "user": os.environ.get(
                "POSTGRES_USER", config.get("postgres", "user", fallback=None)
            ),
            "password": os.environ.get(
                "POSTGRES_PASSWORD",
                config.get("postgres", "password", fallback=None),
            ),
            "database": os.environ.get(
                "POSTGRES_DATABASE",
                config.get("postgres", "database", fallback=None),
            ),
            "workspace": os.environ.get(
                "POSTGRES_WORKSPACE",
                config.get("postgres", "workspace", fallback="default"),
            ),
        }

    def _get_oracle_config():
        return {
            "user": os.environ.get(
                "ORACLE_USER",
                config.get("oracle", "user", fallback=None),
            ),
            "password": os.environ.get(
                "ORACLE_PASSWORD",
                config.get("oracle", "password", fallback=None),
            ),
            "dsn": os.environ.get(
                "ORACLE_DSN",
                config.get("oracle", "dsn", fallback=None),
            ),
            "config_dir": os.environ.get(
                "ORACLE_CONFIG_DIR",
                config.get("oracle", "config_dir", fallback=None),
            ),
            "wallet_location": os.environ.get(
                "ORACLE_WALLET_LOCATION",
                config.get("oracle", "wallet_location", fallback=None),
            ),
            "wallet_password": os.environ.get(
                "ORACLE_WALLET_PASSWORD",
                config.get("oracle", "wallet_password", fallback=None),
            ),
            "workspace": os.environ.get(
                "ORACLE_WORKSPACE",
                config.get("oracle", "workspace", fallback="default"),
            ),
        }

    def _get_tidb_config():
        return {
            "host": os.environ.get(
                "TIDB_HOST",
                config.get("tidb", "host", fallback="localhost"),
            ),
            "port": os.environ.get(
                "TIDB_PORT", config.get("tidb", "port", fallback=4000)
            ),
            "user": os.environ.get(
                "TIDB_USER",
                config.get("tidb", "user", fallback=None),
            ),
            "password": os.environ.get(
                "TIDB_PASSWORD",
                config.get("tidb", "password", fallback=None),
            ),
            "database": os.environ.get(
                "TIDB_DATABASE",
                config.get("tidb", "database", fallback=None),
            ),
            "workspace": os.environ.get(
                "TIDB_WORKSPACE",
                config.get("tidb", "workspace", fallback="default"),
            ),
        }

    # Create the optional API key dependency
    optional_api_key = get_api_key_dependency(api_key)

    # Create working directory if it doesn't exist
    Path(args.working_dir).mkdir(parents=True, exist_ok=True)
    if args.llm_binding == "lollms" or args.embedding_binding == "lollms":
        from lightrag.llm.lollms import lollms_model_complete, lollms_embed
    if args.llm_binding == "ollama" or args.embedding_binding == "ollama":
        from lightrag.llm.ollama import ollama_model_complete, ollama_embed
    if args.llm_binding == "openai" or args.embedding_binding == "openai":
        from lightrag.llm.openai import openai_complete_if_cache, openai_embed
    if args.llm_binding == "azure_openai" or args.embedding_binding == "azure_openai":
        from lightrag.llm.azure_openai import (
            azure_openai_complete_if_cache,
            azure_openai_embed,
        )
    if args.llm_binding_host == "openai-ollama" or args.embedding_binding == "ollama":
        from lightrag.llm.openai import openai_complete_if_cache
        from lightrag.llm.ollama import ollama_embed

    async def openai_alike_model_complete(
        prompt,
        system_prompt=None,
        history_messages=None,
        keyword_extraction=False,
        **kwargs,
    ) -> str:
        keyword_extraction = kwargs.pop("keyword_extraction", None)
        if keyword_extraction:
            kwargs["response_format"] = GPTKeywordExtractionFormat
        if history_messages is None:
            history_messages = []
        return await openai_complete_if_cache(
            args.llm_model,
            prompt,
            system_prompt=system_prompt,
            history_messages=history_messages,
            base_url=args.llm_binding_host,
            api_key=args.llm_binding_api_key,
            **kwargs,
        )

    async def azure_openai_model_complete(
        prompt,
        system_prompt=None,
        history_messages=None,
        keyword_extraction=False,
        **kwargs,
    ) -> str:
        keyword_extraction = kwargs.pop("keyword_extraction", None)
        if keyword_extraction:
            kwargs["response_format"] = GPTKeywordExtractionFormat
        if history_messages is None:
            history_messages = []
        return await azure_openai_complete_if_cache(
            args.llm_model,
            prompt,
            system_prompt=system_prompt,
            history_messages=history_messages,
            base_url=args.llm_binding_host,
            api_key=os.getenv("AZURE_OPENAI_API_KEY"),
            api_version=os.getenv("AZURE_OPENAI_API_VERSION", "2024-08-01-preview"),
            **kwargs,
        )

    embedding_func = EmbeddingFunc(
        embedding_dim=args.embedding_dim,
        max_token_size=args.max_embed_tokens,
        func=lambda texts: lollms_embed(
            texts,
            embed_model=args.embedding_model,
            host=args.embedding_binding_host,
            api_key=args.embedding_binding_api_key,
        )
        if args.embedding_binding == "lollms"
        else ollama_embed(
            texts,
            embed_model=args.embedding_model,
            host=args.embedding_binding_host,
            api_key=args.embedding_binding_api_key,
        )
        if args.embedding_binding == "ollama"
        else azure_openai_embed(
            texts,
            model=args.embedding_model,  # no host is used for openai,
            api_key=args.embedding_binding_api_key,
        )
        if args.embedding_binding == "azure_openai"
        else openai_embed(
            texts,
            model=args.embedding_model,
            base_url=args.embedding_binding_host,
            api_key=args.embedding_binding_api_key,
        ),
    )

    # Initialize RAG
    if args.llm_binding in ["lollms", "ollama", "openai-ollama"]:
        rag = LightRAG(
            working_dir=args.working_dir,
            llm_model_func=lollms_model_complete
            if args.llm_binding == "lollms"
            else ollama_model_complete
            if args.llm_binding == "ollama"
            else openai_alike_model_complete,
            llm_model_name=args.llm_model,
            llm_model_max_async=args.max_async,
            llm_model_max_token_size=args.max_tokens,
            chunk_token_size=int(args.chunk_size),
            chunk_overlap_token_size=int(args.chunk_overlap_size),
            llm_model_kwargs={
                "host": args.llm_binding_host,
                "timeout": args.timeout,
                "options": {"num_ctx": args.max_tokens},
                "api_key": args.llm_binding_api_key,
            }
            if args.llm_binding == "lollms" or args.llm_binding == "ollama"
            else {},
            embedding_func=embedding_func,
            kv_storage=args.kv_storage,
            graph_storage=args.graph_storage,
            vector_storage=args.vector_storage,
            doc_status_storage=args.doc_status_storage,
            vector_db_storage_cls_kwargs={
                "cosine_better_than_threshold": args.cosine_threshold
            },
            enable_llm_cache_for_entity_extract=False,  # set to True for debuging to reduce llm fee
            embedding_cache_config={
                "enabled": True,
                "similarity_threshold": 0.95,
                "use_llm_check": False,
            },
            log_level=args.log_level,
            namespace_prefix=args.namespace_prefix,
        )
    else:
        rag = LightRAG(
            working_dir=args.working_dir,
            llm_model_func=azure_openai_model_complete
            if args.llm_binding == "azure_openai"
            else openai_alike_model_complete,
            chunk_token_size=int(args.chunk_size),
            chunk_overlap_token_size=int(args.chunk_overlap_size),
            llm_model_kwargs={
                "timeout": args.timeout,
            },
            llm_model_name=args.llm_model,
            llm_model_max_async=args.max_async,
            llm_model_max_token_size=args.max_tokens,
            embedding_func=embedding_func,
            kv_storage=args.kv_storage,
            graph_storage=args.graph_storage,
            vector_storage=args.vector_storage,
            doc_status_storage=args.doc_status_storage,
            vector_db_storage_cls_kwargs={
                "cosine_better_than_threshold": args.cosine_threshold
            },
            enable_llm_cache_for_entity_extract=False,  # set to True for debuging to reduce llm fee
            embedding_cache_config={
                "enabled": True,
                "similarity_threshold": 0.95,
                "use_llm_check": False,
            },
            log_level=args.log_level,
            namespace_prefix=args.namespace_prefix,
        )

    # Collect all storage instances
    storage_instances = [
        ("full_docs", rag.full_docs),
        ("text_chunks", rag.text_chunks),
        ("chunk_entity_relation_graph", rag.chunk_entity_relation_graph),
        ("entities_vdb", rag.entities_vdb),
        ("relationships_vdb", rag.relationships_vdb),
        ("chunks_vdb", rag.chunks_vdb),
        ("doc_status", rag.doc_status),
        ("llm_response_cache", rag.llm_response_cache),
    ]

    async def pipeline_enqueue_file(file_path: Path) -> bool:
        """Add a file to the queue for processing

        Args:
            file_path: Path to the saved file
        Returns:
            bool: True if the file was successfully enqueued, False otherwise
        """
        try:
            content = ""
            ext = file_path.suffix.lower()

            file = None
            async with aiofiles.open(file_path, "rb") as f:
                file = await f.read()

            # Process based on file type
            match ext:
                case ".txt" | ".md":
                    content = file.decode("utf-8")
                case ".pdf":
                    if not pm.is_installed("pypdf2"):
                        pm.install("pypdf2")
                    from PyPDF2 import PdfReader
                    from io import BytesIO

                    pdf_file = BytesIO(file)
                    reader = PdfReader(pdf_file)
                    for page in reader.pages:
                        content += page.extract_text() + "\n"
                case ".docx":
                    if not pm.is_installed("docx"):
                        pm.install("docx")
                    from docx import Document
                    from io import BytesIO

                    docx_content = await file.read()
                    docx_file = BytesIO(docx_content)
                    doc = Document(docx_file)
                    content = "\n".join(
                        [paragraph.text for paragraph in doc.paragraphs]
                    )
                case ".pptx":
                    if not pm.is_installed("pptx"):
                        pm.install("pptx")
                    from pptx import Presentation  # type: ignore
                    from io import BytesIO

                    pptx_content = await file.read()
                    pptx_file = BytesIO(pptx_content)
                    prs = Presentation(pptx_file)
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                content += shape.text + "\n"
                case _:
                    logging.error(
                        f"Unsupported file type: {file_path.name} (extension {ext})"
                    )
                    return False

            # Insert into the RAG queue
            if content:
                await rag.apipeline_enqueue_documents(content)
                logging.info(
                    f"Successfully processed and enqueued file: {file_path.name}"
                )
                return True
            else:
                logging.error(
                    f"No content could be extracted from file: {file_path.name}"
                )

        except Exception as e:
            logging.error(
                f"Error processing or enqueueing file {file_path.name}: {str(e)}"
            )
            logging.error(traceback.format_exc())
        finally:
            if file_path.name.startswith(temp_prefix):
                # Clean up the temporary file after indexing
                try:
                    file_path.unlink()
                except Exception as e:
                    logging.error(f"Error deleting file {file_path}: {str(e)}")
        return False

    async def pipeline_index_file(file_path: Path):
        """Index a file

        Args:
            file_path: Path to the saved file
        """
        try:
            if await pipeline_enqueue_file(file_path):
                await rag.apipeline_process_enqueue_documents()

        except Exception as e:
            logging.error(f"Error indexing file {file_path.name}: {str(e)}")
            logging.error(traceback.format_exc())

    async def pipeline_index_files(file_paths: List[Path]):
        """Index multiple files concurrently

        Args:
            file_paths: Paths to the files to index
        """
        if not file_paths:
            return
        try:
            enqueued = False

            if len(file_paths) == 1:
                enqueued = await pipeline_enqueue_file(file_paths[0])
            else:
                tasks = [pipeline_enqueue_file(path) for path in file_paths]
                enqueued = any(await asyncio.gather(*tasks))

            if enqueued:
                await rag.apipeline_process_enqueue_documents()
        except Exception as e:
            logging.error(f"Error indexing files: {str(e)}")
            logging.error(traceback.format_exc())

    async def pipeline_index_texts(texts: List[str]):
        """Index a list of texts

        Args:
            texts: The texts to index
        """
        if not texts:
            return
        await rag.apipeline_enqueue_documents(texts)
        await rag.apipeline_process_enqueue_documents()

    async def save_temp_file(file: UploadFile = File(...)) -> Path:
        """Save the uploaded file to a temporary location

        Args:
            file: The uploaded file

        Returns:
            Path: The path to the saved file
        """
        # Generate unique filename to avoid conflicts
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        unique_filename = f"{temp_prefix}{timestamp}_{file.filename}"

        # Create a temporary file to save the uploaded content
        temp_path = doc_manager.input_dir / "temp" / unique_filename
        temp_path.parent.mkdir(exist_ok=True)

        # Save the file
        with open(temp_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)
        return temp_path

    async def run_scanning_process():
        """Background task to scan and index documents"""
        global scan_progress

        try:
            new_files = doc_manager.scan_directory_for_new_files()
            scan_progress["total_files"] = len(new_files)

            logger.info(f"Found {len(new_files)} new files to index.")
            for file_path in new_files:
                try:
                    with progress_lock:
                        scan_progress["current_file"] = os.path.basename(file_path)

                    await pipeline_index_file(file_path)

                    with progress_lock:
                        scan_progress["indexed_count"] += 1
                        scan_progress["progress"] = (
                            scan_progress["indexed_count"]
                            / scan_progress["total_files"]
                        ) * 100

                except Exception as e:
                    logging.error(f"Error indexing file {file_path}: {str(e)}")

        except Exception as e:
            logging.error(f"Error during scanning process: {str(e)}")
        finally:
            with progress_lock:
                scan_progress["is_scanning"] = False

    @app.post("/documents/scan", dependencies=[Depends(optional_api_key)])
    async def scan_for_new_documents(background_tasks: BackgroundTasks):
        """Trigger the scanning process"""
        global scan_progress

        with progress_lock:
            if scan_progress["is_scanning"]:
                return {"status": "already_scanning"}

            scan_progress["is_scanning"] = True
            scan_progress["indexed_count"] = 0
            scan_progress["progress"] = 0

        # Start the scanning process in the background
        background_tasks.add_task(run_scanning_process)

        return {"status": "scanning_started"}

    @app.get("/documents/scan-progress")
    async def get_scan_progress():
        """Get the current scanning progress"""
        with progress_lock:
            return scan_progress

    @app.post("/documents/upload", dependencies=[Depends(optional_api_key)])
    async def upload_to_input_dir(
        background_tasks: BackgroundTasks, file: UploadFile = File(...)
    ):
        """
        Endpoint for uploading a file to the input directory and indexing it.

        This API endpoint accepts a file through an HTTP POST request, checks if the
        uploaded file is of a supported type, saves it in the specified input directory,
        indexes it for retrieval, and returns a success status with relevant details.

        Parameters:
            background_tasks: FastAPI BackgroundTasks for async processing
            file (UploadFile): The file to be uploaded. It must have an allowed extension as per
                               `doc_manager.supported_extensions`.

        Returns:
            dict: A dictionary containing the upload status ("success"),
                  a message detailing the operation result, and
                  the total number of indexed documents.

        Raises:
            HTTPException: If the file type is not supported, it raises a 400 Bad Request error.
                           If any other exception occurs during the file handling or indexing,
                           it raises a 500 Internal Server Error with details about the exception.
        """
        try:
            if not doc_manager.is_supported_file(file.filename):
                raise HTTPException(
                    status_code=400,
                    detail=f"Unsupported file type. Supported types: {doc_manager.supported_extensions}",
                )

            file_path = doc_manager.input_dir / file.filename
            with open(file_path, "wb") as buffer:
                shutil.copyfileobj(file.file, buffer)

            # Add to background tasks
            background_tasks.add_task(pipeline_index_file, file_path)

            return InsertResponse(
                status="success",
                message=f"File '{file.filename}' uploaded successfully. Processing will continue in background.",
            )
        except Exception as e:
            logging.error(f"Error /documents/upload: {file.filename}: {str(e)}")
            logging.error(traceback.format_exc())
            raise HTTPException(status_code=500, detail=str(e))

    @app.post(
        "/documents/text",
        response_model=InsertResponse,
        dependencies=[Depends(optional_api_key)],
    )
    async def insert_text(
        request: InsertTextRequest, background_tasks: BackgroundTasks
    ):
        """
        Insert text into the Retrieval-Augmented Generation (RAG) system.

        This endpoint allows you to insert text data into the RAG system for later retrieval and use in generating responses.

        Args:
            request (InsertTextRequest): The request body containing the text to be inserted.
            background_tasks: FastAPI BackgroundTasks for async processing

        Returns:
            InsertResponse: A response object containing the status of the operation, a message, and the number of documents inserted.
        """
        try:
            background_tasks.add_task(pipeline_index_texts, [request.text])
            return InsertResponse(
                status="success",
                message="Text successfully received. Processing will continue in background.",
            )
        except Exception as e:
            logging.error(f"Error /documents/text: {str(e)}")
            logging.error(traceback.format_exc())
            raise HTTPException(status_code=500, detail=str(e))

    @app.post(
        "/documents/file",
        response_model=InsertResponse,
        dependencies=[Depends(optional_api_key)],
    )
    async def insert_file(
        background_tasks: BackgroundTasks, file: UploadFile = File(...)
    ):
        """Insert a file directly into the RAG system

        Args:
            background_tasks: FastAPI BackgroundTasks for async processing
            file: Uploaded file

        Returns:
            InsertResponse: Status of the insertion operation

        Raises:
            HTTPException: For unsupported file types or processing errors
        """
        try:
            if not doc_manager.is_supported_file(file.filename):
                raise HTTPException(
                    status_code=400,
                    detail=f"Unsupported file type. Supported types: {doc_manager.supported_extensions}",
                )

            # Create a temporary file to save the uploaded content
            temp_path = save_temp_file(file)

            # Add to background tasks
            background_tasks.add_task(pipeline_index_file, temp_path)

            return InsertResponse(
                status="success",
                message=f"File '{file.filename}' saved successfully. Processing will continue in background.",
            )

        except Exception as e:
            logging.error(f"Error /documents/file: {str(e)}")
            logging.error(traceback.format_exc())
            raise HTTPException(status_code=500, detail=str(e))

    @app.post(
        "/documents/batch",
        response_model=InsertResponse,
        dependencies=[Depends(optional_api_key)],
    )
    async def insert_batch(
        background_tasks: BackgroundTasks, files: List[UploadFile] = File(...)
    ):
        """Process multiple files in batch mode

        Args:
            background_tasks: FastAPI BackgroundTasks for async processing
            files: List of files to process

        Returns:
            InsertResponse: Status of the batch insertion operation

        Raises:
            HTTPException: For processing errors
        """
        try:
            inserted_count = 0
            failed_files = []
            temp_files = []

            for file in files:
                if doc_manager.is_supported_file(file.filename):
                    # Create a temporary file to save the uploaded content
                    temp_files.append(save_temp_file(file))
                    inserted_count += 1
                else:
                    failed_files.append(f"{file.filename} (unsupported type)")

            if temp_files:
                background_tasks.add_task(pipeline_index_files, temp_files)

            # Prepare status message
            if inserted_count == len(files):
                status = "success"
                status_message = f"Successfully inserted all {inserted_count} documents"
            elif inserted_count > 0:
                status = "partial_success"
                status_message = f"Successfully inserted {inserted_count} out of {len(files)} documents"
                if failed_files:
                    status_message += f". Failed files: {', '.join(failed_files)}"
            else:
                status = "failure"
                status_message = "No documents were successfully inserted"
                if failed_files:
                    status_message += f". Failed files: {', '.join(failed_files)}"

            return InsertResponse(status=status, message=status_message)

        except Exception as e:
            logging.error(f"Error /documents/batch: {file.filename}: {str(e)}")
            logging.error(traceback.format_exc())
            raise HTTPException(status_code=500, detail=str(e))

    @app.delete(
        "/documents",
        response_model=InsertResponse,
        dependencies=[Depends(optional_api_key)],
    )
    async def clear_documents():
        """
        Clear all documents from the LightRAG system.

        This endpoint deletes all text chunks, entities vector database, and relationships vector database,
        effectively clearing all documents from the LightRAG system.

        Returns:
            InsertResponse: A response object containing the status, message, and the new document count (0 in this case).
        """
        try:
            rag.text_chunks = []
            rag.entities_vdb = None
            rag.relationships_vdb = None
            return InsertResponse(
                status="success", message="All documents cleared successfully"
            )
        except Exception as e:
            logging.error(f"Error DELETE /documents: {str(e)}")
            logging.error(traceback.format_exc())
            raise HTTPException(status_code=500, detail=str(e))

    @app.post(
        "/query", response_model=QueryResponse, dependencies=[Depends(optional_api_key)]
    )
    async def query_text(request: QueryRequest):
        """
        Handle a POST request at the /query endpoint to process user queries using RAG capabilities.

        Parameters:
            request (QueryRequest): The request object containing the query parameters.
        Returns:
            QueryResponse: A Pydantic model containing the result of the query processing.
                           If a string is returned (e.g., cache hit), it's directly returned.
                           Otherwise, an async generator may be used to build the response.

        Raises:
            HTTPException: Raised when an error occurs during the request handling process,
                           with status code 500 and detail containing the exception message.
        """
        try:
            response = await rag.aquery(
                request.query, param=QueryRequestToQueryParams(request)
            )

            # If response is a string (e.g. cache hit), return directly
            if isinstance(response, str):
                return QueryResponse(response=response)

            # If it's an async generator, decide whether to stream based on stream parameter
            if request.stream or hasattr(response, "__aiter__"):
                result = ""
                async for chunk in response:
                    result += chunk
                return QueryResponse(response=result)
            elif isinstance(response, dict):
                result = json.dumps(response, indent=2)
                return QueryResponse(response=result)
            else:
                return QueryResponse(response=str(response))
        except Exception as e:
            trace_exception(e)
            raise HTTPException(status_code=500, detail=str(e))

    @app.post("/query/stream", dependencies=[Depends(optional_api_key)])
    async def query_text_stream(request: QueryRequest):
        """
        This endpoint performs a retrieval-augmented generation (RAG) query and streams the response.

        Args:
            request (QueryRequest): The request object containing the query parameters.
            optional_api_key (Optional[str], optional): An optional API key for authentication. Defaults to None.

        Returns:
            StreamingResponse: A streaming response containing the RAG query results.
        """
        try:
            params = QueryRequestToQueryParams(request)

            params.stream = True
            response = await rag.aquery(  # Use aquery instead of query, and add await
                request.query, param=params
            )

            from fastapi.responses import StreamingResponse

            async def stream_generator():
                if isinstance(response, str):
                    # If it's a string, send it all at once
                    yield f"{json.dumps({'response': response})}\n"
                else:
                    # If it's an async generator, send chunks one by one
                    try:
                        async for chunk in response:
                            if chunk:  # Only send non-empty content
                                yield f"{json.dumps({'response': chunk})}\n"
                    except Exception as e:
                        logging.error(f"Streaming error: {str(e)}")
                        yield f"{json.dumps({'error': str(e)})}\n"

            return StreamingResponse(
                stream_generator(),
                media_type="application/x-ndjson",
                headers={
                    "Cache-Control": "no-cache",
                    "Connection": "keep-alive",
                    "Content-Type": "application/x-ndjson",
                    "X-Accel-Buffering": "no",  # 确保在Nginx代理时正确处理流式响应
                },
            )
        except Exception as e:
            trace_exception(e)
            raise HTTPException(status_code=500, detail=str(e))

    # query all graph labels
    @app.get("/graph/label/list")
    async def get_graph_labels():
        return await rag.get_graph_labels()

    # query all graph
    @app.get("/graphs")
    async def get_knowledge_graph(label: str):
        return await rag.get_knowledge_graph(nodel_label=label, max_depth=100)

    # Add Ollama API routes
    ollama_api = OllamaAPI(rag, top_k=args.top_k)
    app.include_router(ollama_api.router, prefix="/api")

    @app.get("/documents", dependencies=[Depends(optional_api_key)])
    async def documents():
        """Get current system status"""
        return doc_manager.indexed_files

    @app.get("/health", dependencies=[Depends(optional_api_key)])
    async def get_status():
        """Get current system status"""
        files = doc_manager.scan_directory()
        return {
            "status": "healthy",
            "working_directory": str(args.working_dir),
            "input_directory": str(args.input_dir),
            "indexed_files": [str(f) for f in files],
            "indexed_files_count": len(files),
            "configuration": {
                # LLM configuration binding/host address (if applicable)/model (if applicable)
                "llm_binding": args.llm_binding,
                "llm_binding_host": args.llm_binding_host,
                "llm_model": args.llm_model,
                # embedding model configuration binding/host address (if applicable)/model (if applicable)
                "embedding_binding": args.embedding_binding,
                "embedding_binding_host": args.embedding_binding_host,
                "embedding_model": args.embedding_model,
                "max_tokens": args.max_tokens,
                "kv_storage": args.kv_storage,
                "doc_status_storage": args.doc_status_storage,
                "graph_storage": args.graph_storage,
                "vector_storage": args.vector_storage,
            },
        }

    # Webui mount webui/index.html
    webui_dir = Path(__file__).parent / "webui"
    app.mount(
        "/graph-viewer",
        StaticFiles(directory=webui_dir, html=True),
        name="webui",
    )

    # Serve the static files
    static_dir = Path(__file__).parent / "static"
    static_dir.mkdir(exist_ok=True)
    app.mount("/webui", StaticFiles(directory=static_dir, html=True), name="static")

    return app


def main():
    args = parse_args()
    import uvicorn

    app = create_app(args)
    display_splash_screen(args)
    uvicorn_config = {
        "app": app,
        "host": args.host,
        "port": args.port,
    }
    if args.ssl:
        uvicorn_config.update(
            {
                "ssl_certfile": args.ssl_certfile,
                "ssl_keyfile": args.ssl_keyfile,
            }
        )
    uvicorn.run(**uvicorn_config)


if __name__ == "__main__":
    main()
